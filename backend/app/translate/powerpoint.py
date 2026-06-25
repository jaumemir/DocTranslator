# translate/powerpoint.py
"""
PowerPoint document translation handler (.pptx)

Core functionality:
1. Identify and distinguish titles/subtitles/body/table elements
2. Preserve non-text elements like images and charts
3. Intelligently adjust container size and font to fit translated text
4. Prevent element overlap and boundary overflow
5. Correctly duplicate slides in bilingual mode

Element handling strategies:
- Title/Subtitle: No line breaks, prioritize width expansion or font reduction
- Body text frames: Allow line breaks, can expand height
- Tables: Maintain size, reduce font
- Images/Charts: Preserve completely without processing
"""

import datetime
import logging
import re
import copy
from typing import List, Dict, Any, Optional, Tuple, Set
from dataclasses import dataclass
from threading import Event
from enum import Enum
from pptx import Presentation
from pptx.util import Pt, Emu, Inches
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER
from pptx.shapes.base import BaseShape
from pptx.shapes.placeholder import PlaceholderPicture
from pptx.text.text import TextFrame, _Paragraph, _Run
from pptx.oxml.ns import qn
from pptx.opc.constants import RELATIONSHIP_TYPE as RT
from lxml import etree

from . import to_translate
from . import common

# ==================== Configuration ====================

MAX_CHUNK_SIZE = 2000
MIN_FONT_SIZE = Pt(10)  # Minimum font size
MAX_WIDTH_EXPANSION = 1.3  # Maximum width expansion ratio
MAX_HEIGHT_EXPANSION = 1.5  # Maximum height expansion ratio
FONT_SHRINK_STEP = 0.9  # Font shrink step
MIN_FONT_SCALE = 0.6  # Minimum font scale ratio


class ElementType(Enum):
    """Element type"""
    TITLE = "title"
    SUBTITLE = "subtitle"
    BODY = "body"
    TABLE_CELL = "table_cell"
    TEXT_BOX = "text_box"
    OTHER = "other"


@dataclass
class RunStyle:
    """Run style"""
    font_name: Optional[str] = None
    font_name_ea: Optional[str] = None
    font_size: Optional[int] = None
    bold: Optional[bool] = None
    italic: Optional[bool] = None
    underline: Optional[bool] = None
    color_rgb: Optional[RGBColor] = None


@dataclass
class ShapeGeometry:
    """Shape geometry information"""
    left: int = 0
    top: int = 0
    width: int = 0
    height: int = 0


@dataclass
class TextBlock:
    """Text block"""
    uid: str
    slide_index: int
    shape_id: int
    element_type: ElementType

    # Location information
    location_type: str = "textframe"  # textframe, table_cell
    paragraph_index: int = 0
    cell_row: int = -1
    cell_col: int = -1

    # Text content
    original_text: str = ""
    translated_text: str = ""

    # Style and geometry
    run_style: Optional[RunStyle] = None
    geometry: Optional[ShapeGeometry] = None

    # Status
    complete: bool = False
    skip: bool = False
    count: int = 0

    # Sub-block information
    is_sub: bool = False
    sub_index: int = 0
    sub_total: int = 1
    parent_uid: str = ""


# ==================== Entry Functions ====================

def start(trans: Dict[str, Any]) -> bool:
    """PPT translation entry point"""
    translate_id = trans['id']
    start_time = datetime.datetime.now()

    trans_type = trans.get('type', 'trans_only_inherit')
    only_translation = 'only' in trans_type
    is_bilingual = 'both' in trans_type
    target_lang = trans.get('lang', '英语')

    logging.info(
        f"[Task {translate_id}] PPT translation mode: only={only_translation}, bilingual={is_bilingual}")

    # Open file
    try:
        prs = Presentation(trans['file_path'])
    except Exception as e:
        logging.error(f"[Task {translate_id}] Failed to open PPT: {e}")
        to_translate.error(translate_id, f"Failed to open PPT: {str(e)}")
        return False

    # Get slide dimensions (for boundary detection)
    slide_width = prs.slide_width
    slide_height = prs.slide_height

    # Extract text blocks
    try:
        all_blocks = _extract_all_blocks(prs)
    except Exception as e:
        logging.error(f"[Task {translate_id}] Failed to extract text: {e}")
        to_translate.error(translate_id, f"Failed to extract text: {str(e)}")
        return False

    to_translate_blocks = [b for b in all_blocks if not b.skip]

    if not to_translate_blocks:
        logging.info(f"[Task {translate_id}] No text to translate")
        prs.save(trans['target_file'])
        to_translate.complete(trans, 0, "0 seconds")
        return True

    logging.info(f"[Task {translate_id}] Total {len(all_blocks)} blocks, {len(to_translate_blocks)} need translation")

    # Execute translation
    texts = _blocks_to_api_format(to_translate_blocks)
    event = Event()
    success = to_translate.translate_batch(trans, texts, event)

    if not success:
        return False

    _sync_translation_results(to_translate_blocks, texts)

    # Apply translation
    try:
        if is_bilingual:
            text_count = _apply_bilingual_mode(prs, all_blocks, target_lang,
                                               slide_width, slide_height)
        else:
            text_count = _apply_translation_mode(prs, all_blocks, target_lang,
                                                 slide_width, slide_height)

        prs.save(trans['target_file'])
    except Exception as e:
        logging.error(f"[Task {translate_id}] Save failed: {e}")
        to_translate.error(translate_id, f"Save failed: {str(e)}")
        return False

    end_time = datetime.datetime.now()
    spend_time = common.display_spend(start_time, end_time)
    to_translate.complete(trans, text_count, spend_time)
    return True


# ==================== Text Extraction ====================

def _extract_all_blocks(prs: Presentation) -> List[TextBlock]:
    """Extract all text blocks"""
    blocks = []
    uid_counter = [0]

    def next_uid() -> str:
        uid_counter[0] += 1
        return f"blk_{uid_counter[0]}"

    for slide_idx, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            shape_blocks = _extract_shape_blocks(shape, slide_idx, next_uid)
            blocks.extend(shape_blocks)

    return blocks


def _extract_shape_blocks(shape: BaseShape, slide_idx: int, next_uid) -> List[TextBlock]:
    """Extract text blocks from shape"""
    blocks = []
    shape_id = shape.shape_id

    # 1. Skip pictures
    if _is_picture(shape):
        logging.debug(f"Skip picture: shape_id={shape_id}")
        return blocks

    # 2. Skip charts
    if _is_chart(shape):
        logging.debug(f"Skip chart: shape_id={shape_id}")
        return blocks

    # 3. Skip media/OLE objects
    try:
        skip_types = {
            MSO_SHAPE_TYPE.MEDIA,
            MSO_SHAPE_TYPE.EMBEDDED_OLE_OBJECT,
            MSO_SHAPE_TYPE.LINKED_OLE_OBJECT,
            MSO_SHAPE_TYPE.IGX_GRAPHIC,
        }
        if shape.shape_type in skip_types:
            return blocks
    except:
        pass

    # 4. Recursively process group shapes
    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        try:
            for sub_shape in shape.shapes:
                sub_blocks = _extract_shape_blocks(sub_shape, slide_idx, next_uid)
                blocks.extend(sub_blocks)
        except:
            pass
        return blocks

    # 5. Get shape geometry information
    geometry = _get_shape_geometry(shape)

    # 6. Identify element type
    element_type = _identify_element_type(shape)

    # 7. Tables
    if shape.has_table:
        blocks.extend(_extract_table_blocks(shape.table, slide_idx, shape_id,
                                            geometry, next_uid))
        return blocks

    # 8. Text frames
    if shape.has_text_frame:
        blocks.extend(_extract_textframe_blocks(shape.text_frame, slide_idx, shape_id,
                                                element_type, geometry, next_uid))

    return blocks


def _is_picture(shape: BaseShape) -> bool:
    """Check if shape is a picture"""
    # Check shape type
    try:
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            return True
    except:
        pass

    # Check if it's a picture placeholder
    try:
        if isinstance(shape, PlaceholderPicture):
            return True
    except:
        pass

    # Check if XML contains picture
    try:
        elem = shape._element
        # Check for blip (picture reference)
        blips = elem.findall('.//' + qn('a:blip'))
        if blips:
            # If has blip but no text frame, it's a pure picture
            if not shape.has_text_frame:
                return True
            # If has text frame but empty, also treat as picture
            if shape.has_text_frame and not shape.text_frame.text.strip():
                return True
    except:
        pass

    return False


def _is_chart(shape: BaseShape) -> bool:
    """Check if shape is a chart"""
    try:
        if shape.shape_type == MSO_SHAPE_TYPE.CHART:
            return True
    except:
        pass

    # Check XML
    try:
        elem = shape._element
        charts = elem.findall('.//' + qn('c:chart'))
        if charts:
            return True
    except:
        pass

    try:
        if hasattr(shape, 'chart'):
            return True
    except:
        pass

    return False


def _get_shape_geometry(shape: BaseShape) -> ShapeGeometry:
    """Get shape geometry information"""
    try:
        return ShapeGeometry(
            left=shape.left or 0,
            top=shape.top or 0,
            width=shape.width or 0,
            height=shape.height or 0
        )
    except:
        return ShapeGeometry()


def _identify_element_type(shape: BaseShape) -> ElementType:
    """Identify element type"""
    # Check if it's a placeholder
    try:
        if shape.is_placeholder:
            ph_type = shape.placeholder_format.type

            # Title types
            if ph_type in [PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE]:
                return ElementType.TITLE

            # Subtitle types
            if ph_type in [PP_PLACEHOLDER.SUBTITLE]:
                return ElementType.SUBTITLE

            # Body types
            if ph_type in [PP_PLACEHOLDER.BODY, PP_PLACEHOLDER.OBJECT]:
                return ElementType.BODY
    except:
        pass

    # Non-placeholder text boxes
    if shape.has_text_frame:
        return ElementType.TEXT_BOX

    return ElementType.OTHER


def _extract_textframe_blocks(text_frame: TextFrame, slide_idx: int, shape_id: int,
                              element_type: ElementType, geometry: ShapeGeometry,
                              next_uid) -> List[TextBlock]:
    """Extract paragraphs from text frame"""
    blocks = []

    for para_idx, paragraph in enumerate(text_frame.paragraphs):
        text = _get_paragraph_text(paragraph)

        if not text or not text.strip():
            continue

        run_style = _extract_first_run_style(paragraph)
        should_skip = not _should_translate(text)

        if should_skip:
            blocks.append(TextBlock(
                uid=next_uid(),
                slide_index=slide_idx,
                shape_id=shape_id,
                element_type=element_type,
                location_type='textframe',
                paragraph_index=para_idx,
                original_text=text,
                run_style=run_style,
                geometry=geometry,
                skip=True
            ))
            continue

        # Split into chunks
        if len(text) <= MAX_CHUNK_SIZE:
            blocks.append(TextBlock(
                uid=next_uid(),
                slide_index=slide_idx,
                shape_id=shape_id,
                element_type=element_type,
                location_type='textframe',
                paragraph_index=para_idx,
                original_text=text,
                run_style=run_style,
                geometry=geometry
            ))
        else:
            sub_texts = _split_by_sentences(text, MAX_CHUNK_SIZE)
            parent_uid = next_uid()
            for i, sub_text in enumerate(sub_texts):
                blocks.append(TextBlock(
                    uid=next_uid(),
                    slide_index=slide_idx,
                    shape_id=shape_id,
                    element_type=element_type,
                    location_type='textframe',
                    paragraph_index=para_idx,
                    original_text=sub_text,
                    run_style=run_style,
                    geometry=geometry,
                    is_sub=True,
                    sub_index=i,
                    sub_total=len(sub_texts),
                    parent_uid=parent_uid
                ))

    return blocks


def _extract_table_blocks(table, slide_idx: int, shape_id: int,
                          geometry: ShapeGeometry, next_uid) -> List[TextBlock]:
    """Extract table text"""
    blocks = []
    processed_cells: Set[int] = set()

    for row_idx, row in enumerate(table.rows):
        for col_idx, cell in enumerate(row.cells):
            cell_id = id(cell._tc)
            if cell_id in processed_cells:
                continue
            processed_cells.add(cell_id)

            cell_text = _get_cell_text(cell)

            if not cell_text or not cell_text.strip():
                continue

            run_style = None
            if cell.text_frame.paragraphs:
                run_style = _extract_first_run_style(cell.text_frame.paragraphs[0])

            should_skip = not _should_translate(cell_text)

            if should_skip:
                blocks.append(TextBlock(
                    uid=next_uid(),
                    slide_index=slide_idx,
                    shape_id=shape_id,
                    element_type=ElementType.TABLE_CELL,
                    location_type='table_cell',
                    cell_row=row_idx,
                    cell_col=col_idx,
                    original_text=cell_text,
                    run_style=run_style,
                    geometry=geometry,
                    skip=True
                ))
                continue

            if len(cell_text) <= MAX_CHUNK_SIZE:
                blocks.append(TextBlock(
                    uid=next_uid(),
                    slide_index=slide_idx,
                    shape_id=shape_id,
                    element_type=ElementType.TABLE_CELL,
                    location_type='table_cell',
                    cell_row=row_idx,
                    cell_col=col_idx,
                    original_text=cell_text,
                    run_style=run_style,
                    geometry=geometry
                ))
            else:
                sub_texts = _split_by_sentences(cell_text, MAX_CHUNK_SIZE)
                parent_uid = next_uid()
                for i, sub_text in enumerate(sub_texts):
                    blocks.append(TextBlock(
                        uid=next_uid(),
                        slide_index=slide_idx,
                        shape_id=shape_id,
                        element_type=ElementType.TABLE_CELL,
                        location_type='table_cell',
                        cell_row=row_idx,
                        cell_col=col_idx,
                        original_text=sub_text,
                        run_style=run_style,
                        geometry=geometry,
                        is_sub=True,
                        sub_index=i,
                        sub_total=len(sub_texts),
                        parent_uid=parent_uid
                    ))

    return blocks


def _get_paragraph_text(paragraph: _Paragraph) -> str:
    """Get paragraph text"""
    parts = []
    for run in paragraph.runs:
        if run.text:
            parts.append(run.text)
    return ''.join(parts)


def _get_cell_text(cell) -> str:
    """Get cell text"""
    parts = []
    for para in cell.text_frame.paragraphs:
        para_text = _get_paragraph_text(para)
        if para_text:
            parts.append(para_text)
    return '\n'.join(parts)


def _extract_first_run_style(paragraph: _Paragraph) -> Optional[RunStyle]:
    """Extract style from first valid run in paragraph"""
    for run in paragraph.runs:
        if run.text and run.text.strip():
            return _extract_run_style(run)
    return None


def _extract_run_style(run: _Run) -> RunStyle:
    """Extract run style"""
    style = RunStyle()

    try:
        font = run.font
        style.font_name = font.name

        # East Asian font
        try:
            rPr = run._r.get_or_add_rPr()
            ea = rPr.find(qn('a:ea'))
            if ea is not None:
                style.font_name_ea = ea.get('typeface')
        except:
            pass

        if font.size:
            style.font_size = font.size

        style.bold = font.bold
        style.italic = font.italic
        style.underline = font.underline

        try:
            if font.color and font.color.rgb:
                style.color_rgb = font.color.rgb
        except:
            pass

    except:
        pass

    return style


def _should_translate(text: str) -> bool:
    """Check if text should be translated"""
    text = text.strip()

    if not text:
        return False

    if len(text) < 2:
        return False

    if common.is_all_punc(text):
        return False

    # Pure numbers/symbols
    if re.match(r'^[\d\s\.\-\+\*\/\=\%\(\)\[\]\{\}\#\@\&\|\\\:\;\,\<\>\$\€\¥\£]+$', text):
        return False

    # Page numbers
    if re.match(r'^(第?\s*\d+\s*页?|Page\s*\d+|\d+\s*/\s*\d+)$', text, re.IGNORECASE):
        return False

    # Date time
    if re.match(r'^\d{4}[-/]\d{1,2}[-/]\d{1,2}', text):
        return False

    # URL
    if re.match(r'^https?://', text):
        return False

    # Email
    if re.match(r'^[\w\.-]+@[\w\.-]+\.\w+$', text):
        return False

    return True


def _split_by_sentences(text: str, max_size: int) -> List[str]:
    """Split by sentence boundaries"""
    endings = r'([.!?。！？；;]\s*)'
    parts = re.split(endings, text)

    sentences = []
    i = 0
    while i < len(parts):
        s = parts[i]
        if i + 1 < len(parts) and re.match(endings, parts[i + 1]):
            s += parts[i + 1]
            i += 2
        else:
            i += 1
        if s.strip():
            sentences.append(s)

    chunks = []
    current = ""

    for s in sentences:
        if len(s) > max_size:
            if current:
                chunks.append(current)
                current = ""
            for j in range(0, len(s), max_size):
                chunks.append(s[j:j + max_size])
        elif len(current) + len(s) <= max_size:
            current += s
        else:
            if current:
                chunks.append(current)
            current = s

    if current:
        chunks.append(current)

    return chunks if chunks else [text]


# ==================== Translation Interface ====================

def _blocks_to_api_format(blocks: List[TextBlock]) -> List[Dict]:
    """Convert to API format"""
    return [{
        'text': b.original_text,
        'original': b.original_text,
        'complete': False,
        'count': 0,
        '_uid': b.uid
    } for b in blocks]


def _sync_translation_results(blocks: List[TextBlock], texts: List[Dict]):
    """Sync translation results"""
    uid_map = {b.uid: b for b in blocks}
    for t in texts:
        uid = t.get('_uid')
        if uid and uid in uid_map:
            block = uid_map[uid]
            block.translated_text = t.get('text', block.original_text)
            block.count = t.get('count', 0)
            block.complete = True


# ==================== Translation-Only Mode ====================

def _apply_translation_mode(prs: Presentation, blocks: List[TextBlock],
                            target_lang: str, slide_width: int, slide_height: int) -> int:
    """Apply translation-only mode"""
    text_count = 0

    # Group by slide
    slide_blocks = _group_by_slide(blocks)

    for slide_idx, slide in enumerate(prs.slides):
        if slide_idx not in slide_blocks:
            continue

        sblocks = slide_blocks[slide_idx]
        shape_map = _build_shape_map(slide)

        # Get all shape geometries (for collision detection)
        all_geometries = _get_all_shape_geometries(slide)

        # Group by shape
        shape_blocks = _group_by_shape(sblocks)

        for shape_id, blocks_list in shape_blocks.items():
            if shape_id not in shape_map:
                continue

            shape = shape_map[shape_id]
            count = _apply_to_shape(shape, blocks_list, target_lang,
                                    slide_width, slide_height, all_geometries)
            text_count += count

    return text_count


def _group_by_slide(blocks: List[TextBlock]) -> Dict[int, List[TextBlock]]:
    """Group by slide"""
    result = {}
    for b in blocks:
        if b.slide_index not in result:
            result[b.slide_index] = []
        result[b.slide_index].append(b)
    return result


def _group_by_shape(blocks: List[TextBlock]) -> Dict[int, List[TextBlock]]:
    """Group by shape"""
    result = {}
    for b in blocks:
        if b.shape_id not in result:
            result[b.shape_id] = []
        result[b.shape_id].append(b)
    return result


def _build_shape_map(slide) -> Dict[int, BaseShape]:
    """Build mapping from shape_id to shape"""
    shape_map = {}

    def add_shape(shape):
        shape_map[shape.shape_id] = shape
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            try:
                for sub_shape in shape.shapes:
                    add_shape(sub_shape)
            except:
                pass

    for shape in slide.shapes:
        add_shape(shape)

    return shape_map


def _get_all_shape_geometries(slide) -> List[ShapeGeometry]:
    """Get geometry information of all shapes on slide"""
    geometries = []

    def collect_geometry(shape):
        try:
            geo = _get_shape_geometry(shape)
            if geo.width > 0 and geo.height > 0:
                geometries.append(geo)
        except:
            pass

        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            try:
                for sub_shape in shape.shapes:
                    collect_geometry(sub_shape)
            except:
                pass

    for shape in slide.shapes:
        collect_geometry(shape)

    return geometries


def _apply_to_shape(shape: BaseShape, blocks: List[TextBlock], target_lang: str,
                    slide_width: int, slide_height: int,
                    all_geometries: List[ShapeGeometry]) -> int:
    """Apply translation to shape"""
    if shape.has_table:
        return _apply_to_table(shape.table, blocks, target_lang)
    elif shape.has_text_frame:
        return _apply_to_textframe(shape, blocks, target_lang,
                                   slide_width, slide_height, all_geometries)
    return 0


def _apply_to_textframe(shape: BaseShape, blocks: List[TextBlock], target_lang: str,
                        slide_width: int, slide_height: int,
                        all_geometries: List[ShapeGeometry]) -> int:
    """Apply translation to text frame"""
    text_count = 0
    text_frame = shape.text_frame

    # Group by paragraph
    para_blocks = {}
    for b in blocks:
        if b.location_type != 'textframe':
            continue
        if b.paragraph_index not in para_blocks:
            para_blocks[b.paragraph_index] = []
        para_blocks[b.paragraph_index].append(b)

    # Process each paragraph
    for para_idx, paragraph in enumerate(text_frame.paragraphs):
        if para_idx not in para_blocks:
            continue

        pblocks = sorted(para_blocks[para_idx], key=lambda x: x.sub_index)

        if pblocks[0].skip:
            continue

        # Merge sub-blocks
        if any(b.is_sub for b in pblocks):
            translated = ''.join(b.translated_text or b.original_text for b in pblocks)
            original_text = ''.join(b.original_text for b in pblocks)
        else:
            translated = pblocks[0].translated_text or pblocks[0].original_text
            original_text = pblocks[0].original_text

        text_count += sum(b.count for b in pblocks)
        element_type = pblocks[0].element_type
        run_style = pblocks[0].run_style

        # Decide processing method based on element type
        _replace_paragraph_text_smart(
            paragraph=paragraph,
            shape=shape,
            new_text=translated,
            original_text=original_text,
            element_type=element_type,
            run_style=run_style,
            target_lang=target_lang,
            slide_width=slide_width,
            slide_height=slide_height,
            all_geometries=all_geometries
        )

    return text_count


def _apply_to_table(table, blocks: List[TextBlock], target_lang: str) -> int:
    """Apply translation to table"""
    text_count = 0

    # Group by cell
    cell_blocks = {}
    for b in blocks:
        if b.location_type != 'table_cell':
            continue
        key = (b.cell_row, b.cell_col)
        if key not in cell_blocks:
            cell_blocks[key] = []
        cell_blocks[key].append(b)

    for (row_idx, col_idx), cblocks in cell_blocks.items():
        try:
            cell = table.cell(row_idx, col_idx)
        except:
            continue

        cblocks = sorted(cblocks, key=lambda x: x.sub_index)

        if cblocks[0].skip:
            continue

        if any(b.is_sub for b in cblocks):
            translated = ''.join(b.translated_text or b.original_text for b in cblocks)
            original_text = ''.join(b.original_text for b in cblocks)
        else:
            translated = cblocks[0].translated_text or cblocks[0].original_text
            original_text = cblocks[0].original_text

        text_count += sum(b.count for b in cblocks)

        _replace_cell_text_smart(cell, translated, original_text,
                                 cblocks[0].run_style, target_lang)

    return text_count


def _replace_paragraph_text_smart(paragraph: _Paragraph, shape: BaseShape,
                                  new_text: str, original_text: str,
                                  element_type: ElementType,
                                  run_style: Optional[RunStyle],
                                  target_lang: str,
                                  slide_width: int, slide_height: int,
                                  all_geometries: List[ShapeGeometry]):
    """
    Smart paragraph text replacement
    Adjust container and font based on element type and text length change
    """
    runs = list(paragraph.runs)

    if not runs:
        run = paragraph.add_run()
        run.text = new_text
        _apply_run_style(run, run_style, target_lang)
        return

    # Calculate length change
    original_len = len(original_text)
    new_len = len(new_text)
    length_ratio = new_len / max(original_len, 1)

    # Replace text
    first_run = runs[0]
    first_run.text = new_text
    for run in runs[1:]:
        run.text = ""

    # Ensure font compatibility
    _ensure_font_compatibility(first_run, target_lang, run_style)

    # If text becomes longer, need adjustment
    if length_ratio > 1.1:
        _adjust_for_longer_text(
            paragraph=paragraph,
            shape=shape,
            first_run=first_run,
            length_ratio=length_ratio,
            element_type=element_type,
            run_style=run_style,
            slide_width=slide_width,
            slide_height=slide_height,
            all_geometries=all_geometries
        )


def _adjust_for_longer_text(paragraph: _Paragraph, shape: BaseShape, first_run: _Run,
                            length_ratio: float, element_type: ElementType,
                            run_style: Optional[RunStyle],
                            slide_width: int, slide_height: int,
                            all_geometries: List[ShapeGeometry]):
    """
    Adjust to accommodate longer text
    Strategy:
    1. Title/Subtitle: Prioritize font reduction, then expand width, avoid line breaks
    2. Body/Text box: Prioritize height expansion, then font reduction, allow line breaks
    3. All adjustments must check boundaries and collisions
    """

    if element_type in [ElementType.TITLE, ElementType.SUBTITLE]:
        # Title type: prioritize font reduction
        _adjust_title_element(shape, first_run, length_ratio, run_style,
                              slide_width, slide_height, all_geometries)
    else:
        # Body type: prioritize container expansion
        _adjust_body_element(shape, first_run, length_ratio, run_style,
                             slide_width, slide_height, all_geometries)


def _adjust_title_element(shape: BaseShape, run: _Run, length_ratio: float,
                          run_style: Optional[RunStyle],
                          slide_width: int, slide_height: int,
                          all_geometries: List[ShapeGeometry]):
    """
    Adjust title element
    Strategy: Reduce font > Expand width > Move position
    """
    # 1. First try to reduce font
    font_scale = _calculate_font_scale_for_title(length_ratio)
    if font_scale < 1.0:
        _scale_run_font(run, font_scale, run_style)

    # 2. If still too long, try expanding width
    if length_ratio > 1.5 and font_scale >= MIN_FONT_SCALE:
        try:
            current_width = shape.width
            current_left = shape.left

            # Calculate required new width
            needed_width = int(current_width * min(length_ratio * 0.8, MAX_WIDTH_EXPANSION))

            # Check right boundary
            max_width = slide_width - current_left - Emu(Inches(0.2))  # Leave margin
            new_width = min(needed_width, max_width)

            # Check if it would overlap with other elements
            new_geo = ShapeGeometry(
                left=current_left,
                top=shape.top,
                width=new_width,
                height=shape.height
            )

            if not _would_overlap(new_geo, all_geometries, shape.shape_id, shape):
                shape.width = new_width
            else:
                # If expansion would overlap, further reduce font
                _scale_run_font(run, 0.85, run_style)

        except Exception as e:
            logging.debug(f"Failed to adjust title width: {e}")


def _adjust_body_element(shape: BaseShape, run: _Run, length_ratio: float,
                         run_style: Optional[RunStyle],
                         slide_width: int, slide_height: int,
                         all_geometries: List[ShapeGeometry]):
    """
    Adjust body element
    Strategy: Expand height > Reduce font
    """
    # 1. First try to expand height (allow line breaks)
    if length_ratio > 1.2:
        try:
            current_height = shape.height

            # Calculate required new height
            needed_height = int(current_height * min(length_ratio * 0.9, MAX_HEIGHT_EXPANSION))

            # Check bottom boundary
            max_height = slide_height - shape.top - Emu(Inches(0.2))
            new_height = min(needed_height, max_height)

            # Check if it would overlap with other elements
            new_geo = ShapeGeometry(
                left=shape.left,
                top=shape.top,
                width=shape.width,
                height=new_height
            )

            if not _would_overlap(new_geo, all_geometries, shape.shape_id, shape):
                shape.height = new_height
            else:
                # If expansion would overlap, reduce font
                font_scale = _calculate_font_scale_for_body(length_ratio)
                _scale_run_font(run, font_scale, run_style)

        except Exception as e:
            logging.debug(f"Failed to adjust body height: {e}")
            # On failure, reduce font
            font_scale = _calculate_font_scale_for_body(length_ratio)
            _scale_run_font(run, font_scale, run_style)

    # 2. If ratio is too large, also need to reduce font
    if length_ratio > 1.8:
        font_scale = _calculate_font_scale_for_body(length_ratio)
        _scale_run_font(run, font_scale, run_style)


def _calculate_font_scale_for_title(length_ratio: float) -> float:
    """Calculate font scale ratio for title"""
    if length_ratio <= 1.2:
        return 1.0
    elif length_ratio <= 1.5:
        return 0.90
    elif length_ratio <= 2.0:
        return 0.80
    elif length_ratio <= 2.5:
        return 0.70
    else:
        return MIN_FONT_SCALE


def _calculate_font_scale_for_body(length_ratio: float) -> float:
    """Calculate font scale ratio for body text"""
    if length_ratio <= 1.5:
        return 1.0
    elif length_ratio <= 2.0:
        return 0.95
    elif length_ratio <= 2.5:
        return 0.90
    elif length_ratio <= 3.0:
        return 0.85
    else:
        return 0.75


def _scale_run_font(run: _Run, scale: float, run_style: Optional[RunStyle]):
    """Scale run font size"""
    try:
        # Get current font size
        current_size = run.font.size
        if current_size is None and run_style and run_style.font_size:
            current_size = run_style.font_size
        if current_size is None:
            current_size = Pt(18)  # Default size

        # Calculate new size
        new_size = int(current_size * scale)
        if new_size < MIN_FONT_SIZE:
            new_size = MIN_FONT_SIZE

        run.font.size = new_size
    except Exception as e:
        logging.debug(f"Failed to scale font: {e}")


def _would_overlap(new_geo: ShapeGeometry, all_geometries: List[ShapeGeometry],
                   exclude_shape_id: int, shape: BaseShape) -> bool:
    """Check if new geometry would overlap with other shapes"""
    # Get current shape geometry
    current_geo = _get_shape_geometry(shape)

    for geo in all_geometries:
        # Skip self
        if (geo.left == current_geo.left and geo.top == current_geo.top and
                geo.width == current_geo.width and geo.height == current_geo.height):
            continue

        # Check if overlapping
        if _geometries_overlap(new_geo, geo):
            return True

    return False


def _geometries_overlap(geo1: ShapeGeometry, geo2: ShapeGeometry) -> bool:
    """Check if two geometries overlap"""
    # AABB collision detection
    left1, top1, right1, bottom1 = geo1.left, geo1.top, geo1.left + geo1.width, geo1.top + geo1.height
    left2, top2, right2, bottom2 = geo2.left, geo2.top, geo2.left + geo2.width, geo2.top + geo2.height

    # If not overlapping, return False
    if right1 <= left2 or right2 <= left1:
        return False
    if bottom1 <= top2 or bottom2 <= top1:
        return False

    return True


def _replace_cell_text_smart(cell, new_text: str, original_text: str,
                             run_style: Optional[RunStyle], target_lang: str):
    """Smart cell text replacement"""
    paragraphs = cell.text_frame.paragraphs
    lines = new_text.split('\n')

    # Calculate length change
    length_ratio = len(new_text) / max(len(original_text), 1)
    font_scale = 1.0
    if length_ratio > 1.3:
        font_scale = max(0.7, 1.0 / (length_ratio * 0.8))

    for i, line in enumerate(lines):
        if i < len(paragraphs):
            para = paragraphs[i]
            runs = list(para.runs)

            if runs:
                runs[0].text = line
                for run in runs[1:]:
                    run.text = ""

                if font_scale < 1.0:
                    _scale_run_font(runs[0], font_scale, run_style)
                _ensure_font_compatibility(runs[0], target_lang, run_style)
            else:
                run = para.add_run()
                run.text = line
                if font_scale < 1.0:
                    _scale_run_font(run, font_scale, run_style)
                _apply_run_style(run, run_style, target_lang)
        elif paragraphs:
            # Append to last paragraph
            last_para = paragraphs[-1]
            if last_para.runs:
                last_para.runs[-1].text += '\n' + line


def _apply_run_style(run: _Run, style: Optional[RunStyle], target_lang: str):
    """Apply run style"""
    if not style:
        _set_default_font(run, target_lang)
        return

    try:
        font = run.font

        if style.bold is not None:
            font.bold = style.bold
        if style.italic is not None:
            font.italic = style.italic
        if style.underline is not None:
            font.underline = style.underline
        if style.font_size:
            font.size = style.font_size
        if style.color_rgb:
            font.color.rgb = style.color_rgb

        _set_font_names(run, style, target_lang)

    except:
        pass


def _ensure_font_compatibility(run: _Run, target_lang: str, style: Optional[RunStyle]):
    """Ensure font compatibility with target language"""
    try:
        current_font = run.font.name

        if target_lang in ['中文', '日语', '韩语']:
            if not _is_cjk_font(current_font):
                _set_font_names(run, style, target_lang)
    except:
        pass


def _set_font_names(run: _Run, style: Optional[RunStyle], target_lang: str):
    """Set font names"""
    try:
        # Latin font
        latin_font = 'Arial'
        if style and style.font_name:
            latin_font = style.font_name

        # East Asian font
        ea_font = None
        if style and style.font_name_ea:
            ea_font = style.font_name_ea

        # Set East Asian font based on target language
        if target_lang in ['中文', '日语', '韩语']:
            if not ea_font or not _is_cjk_font(ea_font):
                if target_lang == '中文':
                    ea_font = 'Microsoft YaHei'
                elif target_lang == '日语':
                    ea_font = 'Yu Gothic'
                elif target_lang == '韩语':
                    ea_font = 'Malgun Gothic'

        run.font.name = latin_font

        # Set East Asian font
        if ea_font:
            try:
                rPr = run._r.get_or_add_rPr()

                # Set latin font
                latin = rPr.find(qn('a:latin'))
                if latin is None:
                    latin = etree.SubElement(rPr, qn('a:latin'))
                latin.set('typeface', latin_font)

                # Set ea font
                ea = rPr.find(qn('a:ea'))
                if ea is None:
                    ea = etree.SubElement(rPr, qn('a:ea'))
                ea.set('typeface', ea_font)
            except:
                pass

    except:
        pass


def _set_default_font(run: _Run, target_lang: str):
    """Set default font"""
    try:
        if target_lang in ['中文', '日语', '韩语']:
            run.font.name = 'Microsoft YaHei'
        else:
            run.font.name = 'Arial'
    except:
        pass


def _is_cjk_font(font_name: str) -> bool:
    """Check if it's a CJK font"""
    if not font_name:
        return False

    cjk_keywords = [
        'yahei', '雅黑', 'simsun', '宋体', 'simhei', '黑体',
        'kaiti', '楷体', 'fangsong', '仿宋',
        'gothic', 'mincho', 'meiryo', 'hiragino',
        'malgun', 'batang', 'gulim', 'dotum',
        'noto sans cjk', 'noto serif cjk', 'source han',
        'pingfang', '苹方', 'heiti', 'songti',
        'microsoft jhenghei', '微軟正黑'
    ]

    font_lower = font_name.lower()
    return any(kw in font_lower for kw in cjk_keywords)


# ==================== Bilingual Mode ====================

def _apply_bilingual_mode(prs: Presentation, blocks: List[TextBlock],
                          target_lang: str, slide_width: int, slide_height: int) -> int:
    """
    Bilingual mode: Insert translated slide after each original slide
    """
    text_count = 0

    # Group by slide
    slide_blocks = _group_by_slide(blocks)

    original_count = len(prs.slides)

    # Process from back to front
    for slide_idx in range(original_count - 1, -1, -1):
        try:
            original_slide = prs.slides[slide_idx]

            # Duplicate slide (including all elements and pictures)
            new_slide = _duplicate_slide(prs, original_slide)

            if new_slide is None:
                logging.warning(f"Failed to duplicate slide {slide_idx}")
                continue

            # Move to correct position
            _move_slide(prs, len(prs.slides) - 1, slide_idx + 1)

            # Get new slide and apply translation
            translated_slide = prs.slides[slide_idx + 1]

            if slide_idx in slide_blocks:
                sblocks = slide_blocks[slide_idx]

                # Build position mapping
                shape_mapping = _build_shape_mapping_by_position(original_slide, translated_slide)

                # Statistics
                for b in sblocks:
                    if not b.skip:
                        text_count += b.count

                # Apply translation
                _apply_blocks_with_mapping(translated_slide, sblocks, shape_mapping,
                                           target_lang, slide_width, slide_height)

        except Exception as e:
            logging.error(f"Failed to process slide {slide_idx} in bilingual mode: {e}")
            import traceback
            traceback.print_exc()

    return text_count


def _duplicate_slide(prs: Presentation, source_slide) -> Optional[object]:
    """
    Completely duplicate slide, including all elements and picture relationships

    Key: Correctly handle picture relationship references
    """
    try:
        # Create new slide using source slide layout
        new_slide = prs.slides.add_slide(source_slide.slide_layout)

        # Get source and new slide parts (for handling relationships)
        source_part = source_slide.part
        new_part = new_slide.part

        # Clear shapes auto-generated by layout on new slide
        spTree = new_slide.shapes._spTree
        shapes_to_remove = []

        for child in spTree:
            tag = child.tag
            # Keep nvGrpSpPr and grpSpPr (group properties), delete other shapes
            if tag.endswith('}nvGrpSpPr') or tag.endswith('}grpSpPr'):
                continue
            if (tag.endswith('}sp') or tag.endswith('}pic') or
                    tag.endswith('}graphicFrame') or tag.endswith('}grpSp') or
                    tag.endswith('}cxnSp')):
                shapes_to_remove.append(child)

        for elem in shapes_to_remove:
            spTree.remove(elem)

        # Copy all shapes from source slide and correctly handle relationships
        for shape in source_slide.shapes:
            _copy_shape_with_relationships(shape, spTree, source_part, new_part)

        # Copy background
        _copy_slide_background(source_slide, new_slide)

        return new_slide

    except Exception as e:
        logging.error(f"Failed to duplicate slide: {e}")
        import traceback
        traceback.print_exc()
        return None


def _copy_shape_with_relationships(shape: BaseShape, target_spTree,
                                   source_part, target_part):
    """
    Copy shape and its associated resources (pictures, etc.)

    Key steps:
    1. Deep copy XML element
    2. Find and copy all relationship references (like picture blip)
    3. Update new relationship IDs in copied XML
    """
    try:
        # Deep copy shape's XML element
        new_elem = copy.deepcopy(shape._element)

        # Handle all relationship references in element
        _remap_relationships(new_elem, source_part, target_part)

        # Add to target slide
        target_spTree.append(new_elem)

    except Exception as e:
        logging.warning(f"Failed to copy shape: {e}")


def _remap_relationships(element, source_part, target_part):
    """
    Remap all relationship references in element

    Handled relationship types:
    - a:blip (picture)
    - a:hlinkClick (hyperlink)
    - r:link (external link)
    - c:chart (chart)
    etc.
    """
    # Define namespaces and attributes to process
    EMBED_ATTRS = [
        (qn('r:embed'), RT.IMAGE),  # Picture embed
        (qn('r:link'), None),  # External link
    ]

    LINK_ATTRS = [
        qn('r:id'),  # Generic relationship ID
    ]

    # Collect all rIds that need remapping
    rids_to_remap = set()

    # Find all elements with r:embed attribute (mainly pictures)
    for attr_name, rel_type in EMBED_ATTRS:
        for elem in element.iter():
            rid = elem.get(attr_name)
            if rid:
                rids_to_remap.add((elem, attr_name, rid))

    # Find r:embed in a:blip elements (main picture reference method)
    for blip in element.findall('.//' + qn('a:blip')):
        embed_rid = blip.get(qn('r:embed'))
        if embed_rid:
            rids_to_remap.add((blip, qn('r:embed'), embed_rid))
        link_rid = blip.get(qn('r:link'))
        if link_rid:
            rids_to_remap.add((blip, qn('r:link'), link_rid))

    # Find hyperlinks
    for hlinkClick in element.findall('.//' + qn('a:hlinkClick')):
        rid = hlinkClick.get(qn('r:id'))
        if rid:
            rids_to_remap.add((hlinkClick, qn('r:id'), rid))

    # Find chart references
    for chart in element.findall('.//' + qn('c:chart')):
        rid = chart.get(qn('r:id'))
        if rid:
            rids_to_remap.add((chart, qn('r:id'), rid))

    # Find oleObject
    for oleObj in element.findall('.//' + qn('p:oleObj')):
        rid = oleObj.get(qn('r:id'))
        if rid:
            rids_to_remap.add((oleObj, qn('r:id'), rid))

    # Execute remapping
    for elem, attr_name, old_rid in rids_to_remap:
        try:
            new_rid = _copy_relationship(source_part, target_part, old_rid)
            if new_rid:
                elem.set(attr_name, new_rid)
        except Exception as e:
            logging.debug(f"Failed to remap relationship {old_rid}: {e}")


def _copy_relationship(source_part, target_part, rid: str) -> Optional[str]:
    """
    Copy relationship and its target resource

    Returns new relationship ID
    """
    try:
        # Get source relationship
        source_rels = source_part.rels
        if rid not in source_rels:
            return rid  # Relationship doesn't exist, keep as is

        rel = source_rels[rid]

        # Check if it's an external relationship
        if rel.is_external:
            # External link, directly create new relationship
            new_rid = target_part.relate_to(rel.target_ref, rel.reltype, is_external=True)
            return new_rid

        # Internal relationship, need to copy target resource
        target_resource = rel.target_part

        # Check if target part already has the same resource
        # Judge by comparing blob (binary content)
        existing_rid = _find_existing_resource(target_part, target_resource)
        if existing_rid:
            return existing_rid

        # Copy resource and create new relationship
        new_rid = target_part.relate_to(target_resource, rel.reltype)
        return new_rid

    except Exception as e:
        logging.debug(f"Failed to copy relationship {rid}: {e}")
        return rid  # Keep original rId on failure


def _find_existing_resource(target_part, resource_part) -> Optional[str]:
    """
    Find if the same resource already exists in target part

    Judge by comparing resource blob
    """
    try:
        resource_blob = resource_part.blob

        for rid, rel in target_part.rels.items():
            if rel.is_external:
                continue
            try:
                if rel.target_part.blob == resource_blob:
                    return rid
            except:
                continue

        return None
    except:
        return None


def _copy_slide_background(source_slide, new_slide):
    """Copy slide background"""
    try:
        source_cSld = source_slide._element.find(qn('p:cSld'))
        if source_cSld is not None:
            source_bg = source_cSld.find(qn('p:bg'))
            if source_bg is not None:
                new_cSld = new_slide._element.find(qn('p:cSld'))
                if new_cSld is not None:
                    # Remove existing background
                    old_bg = new_cSld.find(qn('p:bg'))
                    if old_bg is not None:
                        new_cSld.remove(old_bg)

                    # Deep copy background
                    new_bg = copy.deepcopy(source_bg)

                    # Remap relationships in background (like background pictures)
                    _remap_relationships(new_bg, source_slide.part, new_slide.part)

                    # Insert at beginning of cSld
                    new_cSld.insert(0, new_bg)
    except Exception as e:
        logging.debug(f"Failed to copy background: {e}")


def _move_slide(prs: Presentation, from_idx: int, to_idx: int):
    """Move slide position"""
    try:
        sldIdLst = prs.slides._sldIdLst
        slides = list(sldIdLst)

        if from_idx < len(slides):
            slide = slides[from_idx]
            sldIdLst.remove(slide)

            # Adjust target index
            actual_to = to_idx if to_idx < from_idx else to_idx - 1
            actual_to = min(actual_to, len(list(sldIdLst)))

            if actual_to >= len(list(sldIdLst)):
                sldIdLst.append(slide)
            else:
                sldIdLst.insert(actual_to, slide)
    except Exception as e:
        logging.warning(f"Failed to move slide: {e}")


def _build_shape_mapping_by_position(original_slide, new_slide) -> Dict[int, BaseShape]:
    """
    Build shape mapping by position

    Use shape position and size for matching, more reliable than simple ordering
    """
    mapping = {}

    orig_shapes = list(original_slide.shapes)
    new_shapes = list(new_slide.shapes)

    # Find matching new shape for each original shape
    used_new_indices = set()

    for orig_shape in orig_shapes:
        orig_geo = _get_shape_geometry(orig_shape)
        best_match = None
        best_distance = float('inf')
        best_idx = -1

        for idx, new_shape in enumerate(new_shapes):
            if idx in used_new_indices:
                continue

            new_geo = _get_shape_geometry(new_shape)

            # Calculate position and size difference
            distance = (
                    abs(orig_geo.left - new_geo.left) +
                    abs(orig_geo.top - new_geo.top) +
                    abs(orig_geo.width - new_geo.width) +
                    abs(orig_geo.height - new_geo.height)
            )

            # If perfectly matched (or very close)
            if distance < best_distance:
                best_distance = distance
                best_match = new_shape
                best_idx = idx

        if best_match is not None and best_distance < Emu(Inches(0.1)):  # Tolerance 0.1 inch
            mapping[orig_shape.shape_id] = best_match
            used_new_indices.add(best_idx)

            # Recursively process group shapes
            if orig_shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                try:
                    if best_match.shape_type == MSO_SHAPE_TYPE.GROUP:
                        _map_group_by_position(orig_shape, best_match, mapping)
                except:
                    pass

    # For unmatched ones, try matching by order
    unmatched_orig = [s for s in orig_shapes if s.shape_id not in mapping]
    unmatched_new = [s for i, s in enumerate(new_shapes) if i not in used_new_indices]

    for i, orig_shape in enumerate(unmatched_orig):
        if i < len(unmatched_new):
            mapping[orig_shape.shape_id] = unmatched_new[i]

    return mapping


def _map_group_by_position(orig_group, new_group, mapping: Dict[int, BaseShape]):
    """Recursively map group shapes (based on position)"""
    try:
        orig_sub = list(orig_group.shapes)
        new_sub = list(new_group.shapes)

        used_new_indices = set()

        for orig_shape in orig_sub:
            orig_geo = _get_shape_geometry(orig_shape)
            best_match = None
            best_distance = float('inf')
            best_idx = -1

            for idx, new_shape in enumerate(new_sub):
                if idx in used_new_indices:
                    continue

                new_geo = _get_shape_geometry(new_shape)
                distance = (
                        abs(orig_geo.left - new_geo.left) +
                        abs(orig_geo.top - new_geo.top)
                )

                if distance < best_distance:
                    best_distance = distance
                    best_match = new_shape
                    best_idx = idx

            if best_match is not None:
                mapping[orig_shape.shape_id] = best_match
                used_new_indices.add(best_idx)

                if orig_shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                    if best_match.shape_type == MSO_SHAPE_TYPE.GROUP:
                        _map_group_by_position(orig_shape, best_match, mapping)

    except Exception as e:
        logging.debug(f"Failed to map group shapes: {e}")


def _apply_blocks_with_mapping(slide, blocks: List[TextBlock],
                               shape_mapping: Dict[int, BaseShape],
                               target_lang: str,
                               slide_width: int, slide_height: int):
    """Apply translation using mapping"""

    # Group by original shape_id
    shape_blocks = _group_by_shape(blocks)

    # Get geometry information
    all_geometries = _get_all_shape_geometries(slide)

    for orig_shape_id, sblocks in shape_blocks.items():
        if orig_shape_id not in shape_mapping:
            continue

        shape = shape_mapping[orig_shape_id]
        _apply_to_shape(shape, sblocks, target_lang, slide_width, slide_height, all_geometries)
